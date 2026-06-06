"""
zero_file_reader.py — ZeroPointAI File Reader

ZERO_MODULE:    core
ZERO_LAYER:     1
ZERO_ESSENTIAL: false
ZERO_ROLE:      Läser bifogade filer och konverterar till text för LLM
ZERO_DEPENDS:   foundation.py
ZERO_USED_BY:   zero_web_server.py, zero_engine.py

Stöder:
    Text:     .txt, .md, .py, .js, .json, .yaml, .csv, .log, .env
    Dokument: .pdf, .docx, .xlsx, .xls, .pptx
    Bilder:   .jpg, .jpeg, .png, .gif, .webp → base64 för vision
    Arkiv:    .zip (listar innehåll)
"""

from __future__ import annotations

import base64
import io
import logging
import os
from pathlib import Path
from typing import Any, Dict, Optional

log = logging.getLogger(__name__)

# Max storlek per fil (10MB)
MAX_FILE_SIZE = 10 * 1024 * 1024
# Max tecken att skicka till LLM
MAX_TEXT_CHARS = 50_000


def read_file(
    filename: str,
    content:  bytes,
    mime_type: Optional[str] = None,
) -> Dict[str, Any]:
    """
    Läser en bifogad fil och returnerar dict med:
        type:     "text" | "image" | "error"
        text:     extraherad text (för LLM)
        image_b64: base64-kodad bild (för vision-LLM)
        filename: originalfilnamn
        size:     bytes
        summary:  kort beskrivning
    """
    size = len(content)
    ext  = Path(filename).suffix.lower()

    if size > MAX_FILE_SIZE:
        return _error(filename, f"Filen är för stor ({size/1024/1024:.1f}MB, max 10MB)")

    # ── Bilder ────────────────────────────────────────────────────────────────
    if ext in (".jpg", ".jpeg", ".png", ".gif", ".webp"):
        return _read_image(filename, content, mime_type or f"image/{ext[1:]}")

    # ── PDF ───────────────────────────────────────────────────────────────────
    if ext == ".pdf":
        return _read_pdf(filename, content)

    # ── DOCX ──────────────────────────────────────────────────────────────────
    if ext == ".docx":
        return _read_docx(filename, content)

    # ── Excel ─────────────────────────────────────────────────────────────────
    if ext in (".xlsx", ".xls"):
        return _read_excel(filename, content)

    # ── PowerPoint ────────────────────────────────────────────────────────────
    if ext == ".pptx":
        return _read_pptx(filename, content)

    # ── Text / kod ────────────────────────────────────────────────────────────
    if ext in (".txt", ".md", ".py", ".js", ".ts", ".jsx", ".tsx",
               ".json", ".yaml", ".yml", ".csv", ".log", ".env",
               ".sh", ".bash", ".conf", ".ini", ".toml", ".html",
               ".css", ".sql", ".xml", ".rst"):
        return _read_text(filename, content)

    # ── Zip ───────────────────────────────────────────────────────────────────
    if ext == ".zip":
        return _read_zip(filename, content)

    # ── Okänd typ — försök som text ───────────────────────────────────────────
    return _read_text(filename, content, fallback=True)


# ── Readers ───────────────────────────────────────────────────────────────────

def _read_text(filename: str, content: bytes, fallback: bool = False) -> Dict:
    """Läser textfiler och kodfiler."""
    try:
        text = content.decode("utf-8", errors="replace")
        lines = text.splitlines()
        line_count = len(lines)

        if len(text) > MAX_TEXT_CHARS:
            text = text[:MAX_TEXT_CHARS] + f"\n\n...[trunkerad — {len(text)} tecken totalt]"

        ext = Path(filename).suffix.lower()
        lang = _detect_language(ext)

        return {
            "type":     "text",
            "text":     f"**Fil: {filename}** ({line_count} rader)\n\n```{lang}\n{text}\n```",
            "filename": filename,
            "size":     len(content),
            "summary":  f"{filename} ({line_count} rader, {len(content)} bytes)",
        }
    except Exception as e:
        return _error(filename, str(e))


def _read_image(filename: str, content: bytes, mime_type: str) -> Dict:
    """Konverterar bild till base64 för vision-LLM."""
    try:
        b64 = base64.b64encode(content).decode("utf-8")
        # Korrigera mime-type
        if "jpeg" in mime_type or filename.lower().endswith(".jpg"):
            mime_type = "image/jpeg"
        elif "png" in mime_type:
            mime_type = "image/png"
        elif "gif" in mime_type:
            mime_type = "image/gif"
        elif "webp" in mime_type:
            mime_type = "image/webp"

        return {
            "type":       "image",
            "image_b64":  b64,
            "mime_type":  mime_type,
            "text":       f"[Bild bifogad: {filename}]",
            "filename":   filename,
            "size":       len(content),
            "summary":    f"Bild: {filename} ({len(content)/1024:.0f}KB)",
        }
    except Exception as e:
        return _error(filename, str(e))


def _read_pdf(filename: str, content: bytes) -> Dict:
    """Extraherar text från PDF."""
    # Försök pdfplumber → pypdf → fallback
    text = ""
    pages = 0

    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(content)) as pdf:
            pages = len(pdf.pages)
            texts = []
            for i, page in enumerate(pdf.pages[:50]):  # max 50 sidor
                page_text = page.extract_text() or ""
                if page_text.strip():
                    texts.append(f"--- Sida {i+1} ---\n{page_text}")
            text = "\n\n".join(texts)
    except ImportError:
        try:
            import pypdf
            reader = pypdf.PdfReader(io.BytesIO(content))
            pages  = len(reader.pages)
            texts  = []
            for i, page in enumerate(reader.pages[:50]):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    texts.append(f"--- Sida {i+1} ---\n{page_text}")
            text = "\n\n".join(texts)
        except ImportError:
            return _error(filename, "Installera pdfplumber: pip install pdfplumber")
    except Exception as e:
        return _error(filename, f"PDF-fel: {e}")

    if not text.strip():
        return _error(filename, "Kunde inte extrahera text från PDF (kanske skannad bild?)")

    if len(text) > MAX_TEXT_CHARS:
        text = text[:MAX_TEXT_CHARS] + f"\n\n...[trunkerad — {pages} sidor totalt]"

    return {
        "type":     "text",
        "text":     f"**PDF: {filename}** ({pages} sidor)\n\n{text}",
        "filename": filename,
        "size":     len(content),
        "summary":  f"PDF: {filename} ({pages} sidor)",
    }


def _read_docx(filename: str, content: bytes) -> Dict:
    """Extraherar text från Word-dokument."""
    try:
        from docx import Document
        doc   = Document(io.BytesIO(content))
        parts = []

        # Rubriker och stycken
        for para in doc.paragraphs:
            if para.text.strip():
                if para.style.name.startswith("Heading"):
                    parts.append(f"\n## {para.text}")
                else:
                    parts.append(para.text)

        # Tabeller
        for table in doc.tables:
            parts.append("\n[Tabell:]")
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                parts.append(" | ".join(cells))

        text = "\n".join(parts)
        if len(text) > MAX_TEXT_CHARS:
            text = text[:MAX_TEXT_CHARS] + "\n\n...[trunkerad]"

        return {
            "type":     "text",
            "text":     f"**Word-dokument: {filename}**\n\n{text}",
            "filename": filename,
            "size":     len(content),
            "summary":  f"DOCX: {filename} ({len(doc.paragraphs)} stycken)",
        }
    except ImportError:
        return _error(filename, "Installera python-docx: pip install python-docx")
    except Exception as e:
        return _error(filename, f"DOCX-fel: {e}")


def _read_excel(filename: str, content: bytes) -> Dict:
    """Extraherar data från Excel."""
    try:
        import openpyxl
        wb    = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        parts = []

        for sheet_name in wb.sheetnames[:5]:  # max 5 ark
            ws = wb[sheet_name]
            parts.append(f"\n## Ark: {sheet_name}")
            rows = []
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i >= 200:  # max 200 rader per ark
                    parts.append(f"...[{ws.max_row} rader totalt]")
                    break
                cells = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in cells):
                    rows.append(" | ".join(cells))
            parts.extend(rows)

        text = "\n".join(parts)
        if len(text) > MAX_TEXT_CHARS:
            text = text[:MAX_TEXT_CHARS] + "\n\n...[trunkerad]"

        return {
            "type":     "text",
            "text":     f"**Excel: {filename}** ({len(wb.sheetnames)} ark)\n\n{text}",
            "filename": filename,
            "size":     len(content),
            "summary":  f"Excel: {filename} ({len(wb.sheetnames)} ark)",
        }
    except ImportError:
        try:
            import pandas as pd
            dfs   = pd.read_excel(io.BytesIO(content), sheet_name=None)
            parts = []
            for name, df in list(dfs.items())[:5]:
                parts.append(f"\n## Ark: {name}")
                parts.append(df.head(100).to_string())
            text = "\n".join(parts)
            return {
                "type":     "text",
                "text":     f"**Excel: {filename}**\n\n{text[:MAX_TEXT_CHARS]}",
                "filename": filename,
                "size":     len(content),
                "summary":  f"Excel: {filename}",
            }
        except ImportError:
            return _error(filename, "Installera openpyxl: pip install openpyxl")
    except Exception as e:
        return _error(filename, f"Excel-fel: {e}")


def _read_pptx(filename: str, content: bytes) -> Dict:
    """Extraherar text från PowerPoint."""
    try:
        from pptx import Presentation
        prs   = Presentation(io.BytesIO(content))
        parts = []

        for i, slide in enumerate(prs.slides):
            parts.append(f"\n--- Bild {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text)

        text = "\n".join(parts)
        if len(text) > MAX_TEXT_CHARS:
            text = text[:MAX_TEXT_CHARS] + "\n\n...[trunkerad]"

        return {
            "type":     "text",
            "text":     f"**PowerPoint: {filename}** ({len(prs.slides)} bilder)\n\n{text}",
            "filename": filename,
            "size":     len(content),
            "summary":  f"PPTX: {filename} ({len(prs.slides)} bilder)",
        }
    except ImportError:
        return _error(filename, "Installera python-pptx: pip install python-pptx")
    except Exception as e:
        return _error(filename, f"PPTX-fel: {e}")


def _read_zip(filename: str, content: bytes) -> Dict:
    """Listar innehåll i zip-fil."""
    try:
        import zipfile
        with zipfile.ZipFile(io.BytesIO(content)) as zf:
            names = zf.namelist()
            text  = f"**ZIP: {filename}** ({len(names)} filer)\n\n"
            text += "\n".join(f"  {n}" for n in names[:100])
            if len(names) > 100:
                text += f"\n  ...[{len(names)} filer totalt]"

        return {
            "type":     "text",
            "text":     text,
            "filename": filename,
            "size":     len(content),
            "summary":  f"ZIP: {filename} ({len(names)} filer)",
        }
    except Exception as e:
        return _error(filename, f"ZIP-fel: {e}")


def _error(filename: str, msg: str) -> Dict:
    return {
        "type":     "error",
        "text":     f"[Fel vid läsning av {filename}: {msg}]",
        "filename": filename,
        "size":     0,
        "summary":  f"FEL: {filename} — {msg}",
    }


def _detect_language(ext: str) -> str:
    """Mappar filändelse till kodspråk för syntax highlighting."""
    return {
        ".py":   "python",
        ".js":   "javascript",
        ".ts":   "typescript",
        ".jsx":  "jsx",
        ".tsx":  "tsx",
        ".json": "json",
        ".yaml": "yaml",
        ".yml":  "yaml",
        ".sh":   "bash",
        ".bash": "bash",
        ".html": "html",
        ".css":  "css",
        ".sql":  "sql",
        ".md":   "markdown",
        ".xml":  "xml",
        ".toml": "toml",
    }.get(ext, "")


# ── Installationshjälp ────────────────────────────────────────────────────────

def install_dependencies() -> str:
    """Installerar nödvändiga paket."""
    import subprocess
    packages = ["pdfplumber", "python-docx", "openpyxl", "python-pptx"]
    results  = []
    for pkg in packages:
        try:
            r = subprocess.run(
                ["pip", "install", pkg, "--break-system-packages", "-q"],
                capture_output=True, text=True, timeout=60
            )
            results.append(f"{'✓' if r.returncode == 0 else '✗'} {pkg}")
        except Exception as e:
            results.append(f"✗ {pkg}: {e}")
    return "\n".join(results)


if __name__ == "__main__":
    print("Installerar beroenden...")
    print(install_dependencies())
